import streamlit as st

from PyPDF2 import PdfReader

from langchain.text_splitter import CharacterTextSplitter

from langchain.vectorstores import FAISS

from langchain.chains import ConversationalRetrievalChain

from htmlTemplates import css, bot_template, user_template

from langchain_community.chat_models import ChatOllama

from langchain_community.embeddings import OllamaEmbeddings

from docx import Document

from pptx import Presentation

from langdetect import detect

import time

 

def get_pdf_text(pdf_docs):

    text = ""

    for pdf in pdf_docs:

        pdf_reader = PdfReader(pdf)

        for page in pdf_reader.pages:

            text += page.extract_text()

    return text

 

def get_word_text(word_docs):

    text = ""

    for doc in word_docs:

        word_doc = Document(doc)

        for para in word_doc.paragraphs:

            text += para.text + "\n"

    return text

 

def get_ppt_text(ppt_docs):

    text = ""

    for ppt in ppt_docs:

        presentation = Presentation(ppt)

        for slide in presentation.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    text += shape.text + "\n"

    return text

 

def get_text_chunks(text):

    text_splitter = CharacterTextSplitter(

        separator="\n",

        chunk_size=1000,

        chunk_overlap=200,

        length_function=len

    )

    chunks = text_splitter.split_text(text)

    return chunks

 

def get_vectorstore(text_chunks):

    embeddings = OllamaEmbeddings(model="nomic-embed-text")

    vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)

    return vectorstore

 

def get_conversation_chain(vectorstore, language):

    if language == 'fr':

        llm = ChatOllama(model='llama3', language='fr')

    else:

        llm = ChatOllama(model='llama3', language='en')

 

    conversation_chain = ConversationalRetrievalChain.from_llm(

        llm=llm,

        retriever=vectorstore.as_retriever()

    )

    return conversation_chain

 

def handle_userinput(user_question):

    language = detect(user_question)

    st.write(f"Detected language: {language}")  # Debugging line

 

    vectorstore = st.session_state.vectorstore

    conversation_chain = get_conversation_chain(vectorstore, language)

 

    # Start timing

    start_time = time.time()

 

    # Provide an empty chat history for each question

    response = conversation_chain({'question': user_question, 'chat_history': []})

 

    # End timing

    end_time = time.time()

    response_time = end_time - start_time

 

    # Convert response time to minutes

    response_time_minutes = response_time / 60

 

    # Display the user's question and the bot's response

    st.write(user_template.replace("{{MSG}}", user_question), unsafe_allow_html=True)

    st.write(bot_template.replace("{{MSG}}", response['answer']), unsafe_allow_html=True)

    st.write(f"<p>Response Time: {response_time_minutes:.2f} minutes</p>", unsafe_allow_html=True)

 

def main():

    st.set_page_config(page_title="Chat with multiple documents", page_icon=":page_facing_up:")

    st.write(css, unsafe_allow_html=True)

 

    if "vectorstore" not in st.session_state:

        st.session_state.vectorstore = None

 

    st.header("Chat with multiple Documents", ":page_facing_up:")

    user_question = st.text_input("Ask a question about your documents:")

    if user_question:

        handle_userinput(user_question)

 

    with st.sidebar:

        st.subheader("Your documents")

        pdf_docs = st.file_uploader("Upload your PDFs here and click on 'Process'", accept_multiple_files=True)

        word_docs = st.file_uploader("Upload your Word documents here and click on 'Process'", type=["docx"], accept_multiple_files=True)

        ppt_docs = st.file_uploader("Upload your PowerPoint presentations here and click on 'Process'", type=["pptx"], accept_multiple_files=True)

 

        if st.button("Process"):

            with st.spinner("Processing"):

                raw_text = ""

                if pdf_docs:

                    raw_text += get_pdf_text(pdf_docs)

                if word_docs:

                    raw_text += get_word_text(word_docs)

                if ppt_docs:

                    raw_text += get_ppt_text(ppt_docs)

 

                text_chunks = get_text_chunks(raw_text)

                vectorstore = get_vectorstore(text_chunks)

                st.session_state.vectorstore = vectorstore

 

if __name__ == '__main__':

    main()